"""Self-contained PowerPoint export for the UKRI sustainability dashboard.

Mirrors the dashboard's Plotly figures exactly (colors, templates, layouts)
via kaleido for PNG export. The four network figures and the theme matrix
already render with matplotlib in the dashboard, so the export uses the
same matplotlib code paths for those, preserving identical look.
"""
from __future__ import annotations

import io
from collections import Counter
from datetime import date
from typing import Optional

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import networkx as nx
import numpy as np
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

from data_processor import (
    SUSTAINABILITY_THEMES,
    get_or_build_org_index,
    get_or_build_topic_index,
    get_or_build_person_index,
)

# ── Dashboard palettes (mirror app.py exactly) ────────────────────────────
COUNCIL_COLORS = {
    "EPSRC": "#1f77b4", "MRC": "#d62728", "ESRC": "#ff7f0e",
    "NERC": "#2ca02c", "BBSRC": "#9467bd", "AHRC": "#8c564b",
    "STFC": "#e377c2", "GCRF": "#7f7f7f", "Innovate UK": "#17becf",
    "UKRI": "#bcbd22", "UKRI FLF": "#bcbd22", "ISCF": "#aec7e8",
    "Horizon Europe Guarantee": "#ffbb78", "SPF": "#98df8a",
}
THEME_COLORS = px.colors.qualitative.Set2

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ────────────────────────────────────────────────────────────────────────────
# Slide helpers
# ────────────────────────────────────────────────────────────────────────────

def _new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _add_title_slide(prs: Presentation, title: str, subtitle: str, notes: str = "") -> None:
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title is not None:
        slide.shapes.title.text = title
        for p in slide.shapes.title.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(34); r.font.bold = True
                r.font.color.rgb = RGBColor(0x0B, 0x3D, 0x91)
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(2.2), Inches(12), Inches(2))
    tb.text_frame.text = subtitle
    for p in tb.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(20); r.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def _add_section_divider(prs: Presentation, label: str, notes: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(3.0), Inches(12), Inches(1.5))
    tb.text_frame.text = label
    for p in tb.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(40); r.font.bold = True
            r.font.color.rgb = RGBColor(0x0B, 0x3D, 0x91)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def _add_image_slide(prs: Presentation, title: str, png_bytes: bytes, notes: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.6))
    title_tb.text_frame.text = title
    for p in title_tb.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(22); r.font.bold = True
            r.font.color.rgb = RGBColor(0x0B, 0x3D, 0x91)
    # Fit image to a 12.5 × 6.2 in box while preserving its native aspect ratio.
    # python-pptx will respect the native PNG aspect ratio if we pass only width OR height.
    # We measure the actual image dims and pick whichever fits.
    from PIL import Image
    img = Image.open(io.BytesIO(png_bytes))
    iw, ih = img.size
    img_aspect = iw / ih
    frame_w_in, frame_h_in = 12.5, 6.2
    frame_aspect = frame_w_in / frame_h_in
    if img_aspect >= frame_aspect:
        # image is wider — width-limited
        target_w = Inches(frame_w_in)
        target_h = Inches(frame_w_in / img_aspect)
    else:
        # image is taller — height-limited
        target_h = Inches(frame_h_in)
        target_w = Inches(frame_h_in * img_aspect)
    # Center horizontally and vertically within the frame
    left = Inches(0.4) + (Inches(frame_w_in) - target_w) / 2
    top = Inches(0.95) + (Inches(frame_h_in) - target_h) / 2
    slide.shapes.add_picture(io.BytesIO(png_bytes), left, top, width=target_w, height=target_h)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def _add_metrics_slide(prs, title, metrics, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.6))
    title_tb.text_frame.text = title
    for p in title_tb.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(22); r.font.bold = True
            r.font.color.rgb = RGBColor(0x0B, 0x3D, 0x91)
    cols = 2
    rows = (len(metrics) + cols - 1) // cols
    grid_w = Inches(12.5) / cols
    grid_h = Inches(5.8) / max(rows, 1)
    for i, (label, value) in enumerate(metrics):
        r_i, c_i = divmod(i, cols)
        left = Inches(0.4) + c_i * grid_w
        top = Inches(1.1) + r_i * grid_h
        box = slide.shapes.add_textbox(left, top, grid_w - Emu(50000), grid_h - Emu(50000))
        tf = box.text_frame
        tf.word_wrap = True
        tf.text = label
        for p in tf.paragraphs:
            for r in p.runs:
                r.font.size = Pt(13); r.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
        p2 = tf.add_paragraph()
        p2.text = value
        for r in p2.runs:
            r.font.size = Pt(28); r.font.bold = True
            r.font.color.rgb = RGBColor(0x0B, 0x3D, 0x91)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def _add_table_slide(prs, title, headers, rows, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.6))
    title_tb.text_frame.text = title
    for p in title_tb.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(22); r.font.bold = True
            r.font.color.rgb = RGBColor(0x0B, 0x3D, 0x91)
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(0.4), Inches(1.0),
        Inches(12.5), Inches(min(6.0, 0.4 + 0.36 * n_rows)),
    )
    tbl = table_shape.table
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True; r.font.size = Pt(12)
                r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0x0B, 0x3D, 0x91)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(11)
            if i % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xF2, 0xF5, 0xF9)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


# ────────────────────────────────────────────────────────────────────────────
# Figure conversion helpers
# ────────────────────────────────────────────────────────────────────────────

def _plotly_to_png(fig: go.Figure, *, width: int = 1400, height: int = 600,
                   scale: float = 2.0) -> bytes:
    """Render a Plotly figure to PNG via kaleido at scale=2 for crisp output.

    Width/height should match the dashboard's display proportion so the PPT
    figure looks the same shape as on the web app.
    """
    return fig.to_image(format="png", width=width, height=height, scale=scale)


def _mpl_to_png(fig, dpi: int = 220) -> bytes:
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=dpi, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return buf.getvalue()


def _placeholder_png(text: str) -> bytes:
    fig, ax = plt.subplots(figsize=(10, 5))
    ax.axis("off")
    ax.text(0.5, 0.5, text, ha="center", va="center", fontsize=16, color="#555")
    return _mpl_to_png(fig)


# ────────────────────────────────────────────────────────────────────────────
# Tab 1 — Overview (mirrors app.py Plotly definitions exactly)
# ────────────────────────────────────────────────────────────────────────────

def _fig_annual_funding(df: pd.DataFrame) -> bytes:
    yr_df = (
        df.groupby(["start_year", "funder"])["fund_value"].sum().reset_index()
          .rename(columns={"fund_value": "funding", "start_year": "year", "funder": "council"})
    )
    yr_df["funding_B"] = yr_df["funding"] / 1e9
    fig = px.bar(
        yr_df, x="year", y="funding_B", color="council",
        color_discrete_map=COUNCIL_COLORS,
        labels={"funding_B": "Funding (£B)", "year": "Start Year", "council": "Council"},
        template="plotly_white",
    )
    fig.update_layout(legend_title_text="Council", height=380, margin=dict(t=10))
    return _plotly_to_png(fig, width=1400, height=560)


def _fig_funding_share_donut(df: pd.DataFrame) -> bytes:
    council_df = (
        df.groupby("funder")["fund_value"].sum().reset_index()
          .rename(columns={"fund_value": "funding", "funder": "council"})
          .sort_values("funding", ascending=False)
    )
    fig = px.pie(
        council_df, names="council", values="funding", hole=0.45,
        color="council", color_discrete_map=COUNCIL_COLORS, template="plotly_white",
    )
    fig.update_traces(textposition="inside", textinfo="percent+label")
    fig.update_layout(showlegend=False, height=380, margin=dict(t=10))
    return _plotly_to_png(fig, width=1100, height=720)


def _fig_sustainability_trend(df: pd.DataFrame) -> bytes:
    sus_yr = df.groupby(["start_year", "is_sustainability"])["fund_value"].sum().reset_index()
    sus_yr["type"] = sus_yr["is_sustainability"].map({True: "Sustainability", False: "Other UKRI"})
    sus_yr["funding_M"] = sus_yr["fund_value"] / 1e6
    fig = px.area(
        sus_yr, x="start_year", y="funding_M", color="type",
        color_discrete_map={"Sustainability": "#2ca02c", "Other UKRI": "#aec7e8"},
        labels={"funding_M": "Funding (£M)", "start_year": "Start Year", "type": "Category"},
        template="plotly_white",
    )
    fig.update_layout(height=320, margin=dict(t=10), legend_title_text="")
    return _plotly_to_png(fig, width=1400, height=520)


def _fig_project_count_by_year(df: pd.DataFrame) -> bytes:
    cnt_df = df.groupby(["start_year", "status"]).size().reset_index(name="count")
    fig = px.bar(
        cnt_df, x="start_year", y="count", color="status",
        color_discrete_map={"Active": "#2ca02c", "Closed": "#888"},
        labels={"count": "Projects", "start_year": "Start Year", "status": "Status"},
        template="plotly_white", barmode="stack",
    )
    fig.update_layout(height=320, margin=dict(t=10), legend_title_text="")
    return _plotly_to_png(fig, width=1400, height=520)


def _fig_grant_size_box(df: pd.DataFrame) -> bytes:
    box_df = df[df["fund_value"].notna() & (df["fund_value"] > 0)].copy()
    box_df["funding_M"] = box_df["fund_value"] / 1e6
    fig = px.box(
        box_df, x="funder", y="funding_M", color="funder",
        color_discrete_map=COUNCIL_COLORS, points=False,
        labels={"funding_M": "Grant Size (£M)", "funder": "Council"},
        template="plotly_white", log_y=True,
    )
    fig.update_layout(showlegend=False, height=340, margin=dict(t=10))
    return _plotly_to_png(fig, width=1400, height=560)


# ────────────────────────────────────────────────────────────────────────────
# Tab 2 — Geographic
# ────────────────────────────────────────────────────────────────────────────

def _fig_region_funding(df: pd.DataFrame) -> bytes:
    reg_df = (
        df.groupby("region").agg(
            funding=("fund_value", "sum"),
            projects=("project_ref", "count"),
            sus_projects=("is_sustainability", "sum"),
        ).reset_index().sort_values("funding", ascending=True)
    )
    reg_df["funding_M"] = reg_df["funding"] / 1e6
    fig = go.Figure()
    fig.add_trace(go.Bar(
        y=reg_df["region"], x=reg_df["funding_M"], orientation="h",
        marker_color="#1f77b4", name="Total Funding",
        text=[f"£{v:.0f}M" for v in reg_df["funding_M"]], textposition="outside",
    ))
    fig.update_layout(
        xaxis_title="Funding (£M)", yaxis_title="",
        template="plotly_white", height=500, margin=dict(t=10, r=80),
    )
    return _plotly_to_png(fig, width=1400, height=700)


def _fig_projects_by_region(df: pd.DataFrame) -> bytes:
    reg_df = (
        df.groupby("region").agg(
            projects=("project_ref", "count"),
            sus_projects=("is_sustainability", "sum"),
        ).reset_index()
    )
    reg_df["sus_pct"] = (100 * reg_df["sus_projects"] / reg_df["projects"]).round(1)
    fig = px.bar(
        reg_df.sort_values("projects"),
        y="region", x="projects", orientation="h",
        color="sus_pct", color_continuous_scale="Greens",
        labels={"projects": "Projects", "region": "", "sus_pct": "% Sustainability"},
        template="plotly_white",
    )
    fig.update_layout(height=500, margin=dict(t=10), coloraxis_colorbar_title="% Sust.")
    return _plotly_to_png(fig, width=1400, height=700)


def _fig_nation_funding(df: pd.DataFrame) -> bytes:
    nation_df = df.groupby("nation").agg(
        funding=("fund_value", "sum"), projects=("project_ref", "count")
    ).reset_index()
    sus_nation = (
        df[df["is_sustainability"]].groupby("nation")["fund_value"].sum().reset_index()
          .rename(columns={"fund_value": "sus_funding"})
    )
    nation_df = nation_df.merge(sus_nation, on="nation", how="left").fillna(0)
    nation_df["funding_M"] = nation_df["funding"] / 1e6
    nation_df["sus_M"] = nation_df["sus_funding"] / 1e6
    nation_df["other_M"] = nation_df["funding_M"] - nation_df["sus_M"]
    fig = go.Figure()
    fig.add_trace(go.Bar(name="Sustainability", x=nation_df["nation"], y=nation_df["sus_M"],
                         marker_color="#2ca02c"))
    fig.add_trace(go.Bar(name="Other UKRI", x=nation_df["nation"], y=nation_df["other_M"],
                         marker_color="#aec7e8"))
    fig.update_layout(
        barmode="stack", yaxis_title="Funding (£M)", template="plotly_white",
        height=340, margin=dict(t=10), legend_title_text="",
    )
    return _plotly_to_png(fig, width=1300, height=520)


def _fig_regional_heatmap(df: pd.DataFrame) -> bytes:
    heat_df = df.groupby(["region", "start_year"])["fund_value"].sum().reset_index()
    heat_pivot = heat_df.pivot(index="region", columns="start_year", values="fund_value").fillna(0) / 1e6
    fig = px.imshow(
        heat_pivot, aspect="auto", color_continuous_scale="YlGn",
        labels=dict(color="£M"), template="plotly_white",
    )
    fig.update_layout(height=400, margin=dict(t=10))
    return _plotly_to_png(fig, width=1400, height=560)


# ────────────────────────────────────────────────────────────────────────────
# Tab 3 — Institutions
# ────────────────────────────────────────────────────────────────────────────

def _fig_top_institutions(df: pd.DataFrame) -> bytes:
    inst_df = (
        df.groupby("institution").agg(
            funding=("fund_value", "sum"),
            projects=("project_ref", "count"),
            sus_projects=("is_sustainability", "sum"),
        ).reset_index().sort_values("funding", ascending=False)
    )
    inst_df["funding_M"] = inst_df["funding"] / 1e6
    inst_df["sus_pct"] = (100 * inst_df["sus_projects"] / inst_df["projects"]).round(1)
    top20 = inst_df.head(20).sort_values("funding_M")
    fig = go.Figure()
    fig.add_trace(go.Bar(
        y=top20["institution"], x=top20["funding_M"], orientation="h",
        marker=dict(color=top20["sus_pct"], colorscale="Greens",
                    colorbar=dict(title="% Sustainability")),
        text=[f"£{v:.0f}M" for v in top20["funding_M"]], textposition="outside",
    ))
    fig.update_layout(
        xaxis_title="Total Funding (£M)", yaxis_title="",
        template="plotly_white", height=560, margin=dict(t=10, r=80),
    )
    return _plotly_to_png(fig, width=1400, height=780)


def _fig_institution_bubble(df: pd.DataFrame) -> bytes:
    inst_df = (
        df.groupby("institution").agg(
            funding=("fund_value", "sum"), projects=("project_ref", "count"),
            sus_projects=("is_sustainability", "sum"),
            total_pubs=("publication_count", "sum"),
        ).reset_index().sort_values("funding", ascending=False)
    )
    inst_df["funding_M"] = inst_df["funding"] / 1e6
    inst_df["sus_pct"] = (100 * inst_df["sus_projects"] / inst_df["projects"]).round(1)
    top50 = inst_df[inst_df["projects"] >= 3].head(50)
    fig = px.scatter(
        top50, x="projects", y="funding_M", size="total_pubs",
        color="sus_pct", color_continuous_scale="Greens",
        hover_name="institution",
        labels={
            "projects": "Number of Projects", "funding_M": "Total Funding (£M)",
            "total_pubs": "Publications", "sus_pct": "% Sustainability",
        },
        template="plotly_white", size_max=40,
    )
    fig.update_layout(height=420, margin=dict(t=10))
    return _plotly_to_png(fig, width=1400, height=620)


def _fig_department_funding(df: pd.DataFrame) -> bytes:
    dept_df = (
        df[df["department"].notna() & (df["department"].str.upper() != "UNLISTED")]
          .groupby("department")["fund_value"].sum().reset_index()
          .sort_values("fund_value", ascending=False).head(15)
    )
    dept_df["funding_M"] = dept_df["fund_value"] / 1e6
    fig = px.bar(
        dept_df.sort_values("funding_M"), y="department", x="funding_M",
        orientation="h", color="funding_M", color_continuous_scale="Blues",
        labels={"funding_M": "Funding (£M)", "department": ""},
        template="plotly_white",
    )
    fig.update_layout(showlegend=False, height=420, margin=dict(t=10), coloraxis_showscale=False)
    return _plotly_to_png(fig, width=1400, height=620)


def _fig_institution_council_matrix(df: pd.DataFrame) -> bytes:
    inst_df = (
        df.groupby("institution").agg(funding=("fund_value", "sum"))
          .reset_index().sort_values("funding", ascending=False)
    )
    inst_council = df.groupby(["institution", "funder"])["fund_value"].sum().reset_index()
    top15_inst = inst_df.head(15)["institution"].tolist()
    ic_pivot = (
        inst_council[inst_council["institution"].isin(top15_inst)]
          .pivot(index="institution", columns="funder", values="fund_value")
          .fillna(0)
    ) / 1e6
    fig = px.imshow(
        ic_pivot, aspect="auto", color_continuous_scale="Blues",
        labels=dict(color="£M"), template="plotly_white", text_auto=".0f",
    )
    fig.update_layout(height=420, margin=dict(t=10))
    return _plotly_to_png(fig, width=1400, height=620)


# ────────────────────────────────────────────────────────────────────────────
# Tab 4 — Sustainability themes
# ────────────────────────────────────────────────────────────────────────────

def _build_theme_df(df: pd.DataFrame) -> pd.DataFrame:
    theme_counts = {}
    theme_funding = {}
    for theme in SUSTAINABILITY_THEMES:
        col_key = "theme_" + theme.lower().replace(" ", "_").replace("&", "and")
        if col_key in df.columns:
            mask = df[col_key]
            theme_counts[theme] = mask.sum()
            theme_funding[theme] = df.loc[mask, "fund_value"].sum()
    th_df = pd.DataFrame({
        "theme": list(theme_counts.keys()),
        "projects": list(theme_counts.values()),
        "funding_M": [v / 1e6 for v in theme_funding.values()],
    }).sort_values("projects", ascending=True)
    return th_df


def _fig_themes_bar(df: pd.DataFrame) -> bytes:
    th_df = _build_theme_df(df)
    fig = px.bar(
        th_df, y="theme", x="projects", orientation="h",
        color="funding_M", color_continuous_scale="Greens",
        labels={"projects": "Projects", "theme": "", "funding_M": "Funding (£M)"},
        template="plotly_white",
    )
    fig.update_layout(height=400, margin=dict(t=10), coloraxis_colorbar_title="£M")
    return _plotly_to_png(fig, width=1400, height=600)


def _fig_theme_funding_pie(df: pd.DataFrame) -> bytes:
    th_df = _build_theme_df(df)
    fig = px.pie(
        th_df.sort_values("funding_M", ascending=False),
        names="theme", values="funding_M", hole=0.4,
        color_discrete_sequence=THEME_COLORS, template="plotly_white",
    )
    fig.update_traces(textposition="inside", textinfo="percent+label")
    fig.update_layout(showlegend=False, height=400, margin=dict(t=10))
    return _plotly_to_png(fig, width=1100, height=720)


def _fig_theme_evolution(df: pd.DataFrame) -> bytes:
    rows = []
    for theme in SUSTAINABILITY_THEMES:
        col_key = "theme_" + theme.lower().replace(" ", "_").replace("&", "and")
        if col_key in df.columns:
            sub = df[df[col_key]].groupby("start_year").size().reset_index(name="count")
            sub["theme"] = theme
            rows.append(sub)
    if not rows:
        return _placeholder_png("No theme-year data")
    theme_yr_df = pd.concat(rows)
    fig = px.line(
        theme_yr_df, x="start_year", y="count", color="theme",
        color_discrete_sequence=THEME_COLORS, markers=True,
        labels={"count": "Projects", "start_year": "Start Year", "theme": "Theme"},
        template="plotly_white",
    )
    fig.update_layout(height=380, margin=dict(t=10), legend_title_text="")
    return _plotly_to_png(fig, width=1400, height=580)


def _fig_theme_council_heatmap(df: pd.DataFrame) -> bytes:
    rows = []
    for theme in SUSTAINABILITY_THEMES:
        col_key = "theme_" + theme.lower().replace(" ", "_").replace("&", "and")
        if col_key in df.columns:
            sub = df[df[col_key]].groupby("funder").size().reset_index(name="count")
            sub["theme"] = theme
            rows.append(sub)
    if not rows:
        return _placeholder_png("No theme-council data")
    tc_df = pd.concat(rows)
    tc_pivot = tc_df.pivot(index="theme", columns="funder", values="count").fillna(0)
    fig = px.imshow(
        tc_pivot, aspect="auto", color_continuous_scale="YlGn",
        labels=dict(color="Projects"), template="plotly_white", text_auto="d",
    )
    fig.update_layout(height=380, margin=dict(t=10))
    return _plotly_to_png(fig, width=1400, height=580)


def _fig_theme_region_heatmap(df: pd.DataFrame) -> bytes:
    rows = []
    for theme in SUSTAINABILITY_THEMES:
        col_key = "theme_" + theme.lower().replace(" ", "_").replace("&", "and")
        if col_key in df.columns:
            sub = df[df[col_key]].groupby("region").size().reset_index(name="count")
            sub["theme"] = theme
            rows.append(sub)
    if not rows:
        return _placeholder_png("No theme-region data")
    tr_df = pd.concat(rows)
    tr_pivot = tr_df.pivot(index="theme", columns="region", values="count").fillna(0)
    fig = px.imshow(
        tr_pivot, aspect="auto", color_continuous_scale="Greens",
        labels=dict(color="Projects"), template="plotly_white", text_auto="d",
    )
    fig.update_layout(height=380, margin=dict(t=10))
    return _plotly_to_png(fig, width=1400, height=580)


def _fig_research_topics(df: pd.DataFrame) -> bytes:
    sus = df[df["is_sustainability"]]
    topics_text = " ".join(sus.get("research_topics", pd.Series(dtype=str)).fillna("").tolist())
    items = [t.strip() for t in topics_text.split(";") if t.strip() and t.strip() != "Unclassified"]
    counts = Counter(items).most_common(25)
    if not counts:
        return _placeholder_png("No research-topic data")
    td_df = pd.DataFrame(counts, columns=["topic", "count"]).sort_values("count")
    fig = px.bar(
        td_df, y="topic", x="count", orientation="h",
        color="count", color_continuous_scale="Greens",
        labels={"count": "Occurrences", "topic": ""},
        template="plotly_white",
    )
    fig.update_layout(showlegend=False, height=500, margin=dict(t=10), coloraxis_showscale=False)
    return _plotly_to_png(fig, width=1400, height=720)


# ────────────────────────────────────────────────────────────────────────────
# Tab 5 — Collaboration (non-network charts)
# ────────────────────────────────────────────────────────────────────────────

def _fig_collab_histogram(df: pd.DataFrame) -> bytes:
    collab_df = df[df["collab_count"] > 0].copy()
    if collab_df.empty:
        return _placeholder_png("No collaboration data")
    fig = px.histogram(
        collab_df, x="collab_count", nbins=40,
        color_discrete_sequence=["#17becf"],
        labels={"collab_count": "Number of Partner Organisations", "count": "Projects"},
        template="plotly_white",
    )
    fig.update_layout(height=360, margin=dict(t=10), yaxis_title="Projects")
    return _plotly_to_png(fig, width=1400, height=560)


def _fig_collab_box_by_council(df: pd.DataFrame) -> bytes:
    sub = df[df["collab_count"] > 0]
    if sub.empty:
        return _placeholder_png("No collaboration data")
    fig = px.box(
        sub, x="funder", y="collab_count", color="funder",
        color_discrete_map=COUNCIL_COLORS, points=False,
        labels={"collab_count": "Partner Organisations", "funder": "Council"},
        template="plotly_white",
    )
    fig.update_layout(showlegend=False, height=360, margin=dict(t=10))
    return _plotly_to_png(fig, width=1400, height=560)


def _fig_collab_region_bar(df: pd.DataFrame) -> bytes:
    collab_reg = df.groupby(["region", "is_sustainability"])["collab_count"].mean().reset_index()
    collab_reg["type"] = collab_reg["is_sustainability"].map({True: "Sustainability", False: "Other"})
    fig = px.bar(
        collab_reg, x="region", y="collab_count", color="type", barmode="group",
        color_discrete_map={"Sustainability": "#2ca02c", "Other": "#aec7e8"},
        labels={"collab_count": "Avg. Partner Organisations", "region": "Region", "type": ""},
        template="plotly_white",
    )
    fig.update_xaxes(tickangle=30)
    fig.update_layout(height=360, margin=dict(t=10, b=80), legend_title_text="")
    return _plotly_to_png(fig, width=1500, height=580)


def _fig_collab_vs_grant_scatter(df: pd.DataFrame) -> bytes:
    scat_df = df[(df["collab_count"] > 0) & df["fund_value"].notna()].copy()
    if scat_df.empty:
        return _placeholder_png("No collab/grant data")
    scat_df["funding_M"] = scat_df["fund_value"] / 1e6
    fig = px.scatter(
        scat_df, x="collab_count", y="funding_M",
        color="funder", color_discrete_map=COUNCIL_COLORS,
        opacity=0.6, hover_name="title",
        labels={
            "collab_count": "Number of Collaborating Organisations",
            "funding_M": "Grant Size (£M)", "funder": "Council",
        },
        template="plotly_white",
        trendline="ols", trendline_scope="overall", trendline_color_override="red",
    )
    fig.update_layout(height=420, margin=dict(t=10))
    return _plotly_to_png(fig, width=1400, height=620)


def _fig_top_collaborative_institutions(df: pd.DataFrame) -> bytes:
    inst_collab = (
        df[df["collab_count"] > 0].groupby("institution")
          .agg(avg_collab=("collab_count", "mean"), projects=("project_ref", "count"))
          .reset_index().query("projects >= 3")
          .sort_values("avg_collab", ascending=False).head(20)
          .sort_values("avg_collab")
    )
    if inst_collab.empty:
        return _placeholder_png("No collab institution data")
    fig = px.bar(
        inst_collab, y="institution", x="avg_collab", orientation="h",
        color="projects", color_continuous_scale="Blues",
        labels={"avg_collab": "Avg. Partner Organisations", "institution": "", "projects": "Projects"},
        template="plotly_white",
    )
    fig.update_layout(height=500, margin=dict(t=10), coloraxis_colorbar_title="Projects")
    return _plotly_to_png(fig, width=1400, height=720)


# ────────────────────────────────────────────────────────────────────────────
# Tab 6 — Research impact
# ────────────────────────────────────────────────────────────────────────────

def _fig_pubs_per_project_by_council(df: pd.DataFrame) -> bytes:
    pub_council = (
        df.groupby("funder").agg(
            total_pubs=("publication_count", "sum"),
            projects=("project_ref", "count"),
        ).reset_index()
    )
    pub_council["pubs_per_project"] = pub_council["total_pubs"] / pub_council["projects"]
    pub_council = pub_council.sort_values("pubs_per_project", ascending=True)
    fig = px.bar(
        pub_council, y="funder", x="pubs_per_project", orientation="h",
        color="total_pubs", color_continuous_scale="Blues",
        labels={"pubs_per_project": "Publications per Project", "funder": "", "total_pubs": "Total Pubs"},
        template="plotly_white",
    )
    fig.update_layout(height=380, margin=dict(t=10), coloraxis_colorbar_title="Total Pubs")
    return _plotly_to_png(fig, width=1400, height=580)


def _fig_pubs_histogram(df: pd.DataFrame) -> bytes:
    pub_sus = df.copy()
    pub_sus["type"] = pub_sus["is_sustainability"].map({True: "Sustainability", False: "Other UKRI"})
    fig = px.histogram(
        pub_sus[pub_sus["publication_count"] > 0], x="publication_count",
        color="type", nbins=50, barmode="overlay", opacity=0.7,
        color_discrete_map={"Sustainability": "#2ca02c", "Other UKRI": "#aec7e8"},
        labels={"publication_count": "Publications per Project", "count": "Projects"},
        template="plotly_white", log_y=True,
    )
    fig.update_layout(height=380, margin=dict(t=10), legend_title_text="", yaxis_title="Projects (log)")
    return _plotly_to_png(fig, width=1400, height=580)


def _fig_pubs_vs_grant_scatter(df: pd.DataFrame) -> bytes:
    pub_scat = df[(df["publication_count"] > 0) & df["fund_value"].notna() & (df["fund_value"] > 0)].copy()
    if pub_scat.empty:
        return _placeholder_png("No publication/grant data")
    pub_scat["funding_M"] = pub_scat["fund_value"] / 1e6
    fig = px.scatter(
        pub_scat, x="funding_M", y="publication_count",
        color="funder", color_discrete_map=COUNCIL_COLORS, opacity=0.6, hover_name="title",
        log_x=True, log_y=True,
        labels={
            "funding_M": "Grant Size (£M, log)",
            "publication_count": "Publications (log)",
            "funder": "Council",
        },
        template="plotly_white",
        trendline="ols", trendline_scope="overall", trendline_color_override="black",
    )
    fig.update_layout(height=420, margin=dict(t=10))
    return _plotly_to_png(fig, width=1400, height=620)


def _fig_publication_output_by_region(df: pd.DataFrame) -> bytes:
    pub_reg = (
        df.groupby("region").agg(
            total_pubs=("publication_count", "sum"),
            projects=("project_ref", "count"),
        ).reset_index()
    )
    pub_reg["pubs_per_project"] = (pub_reg["total_pubs"] / pub_reg["projects"]).round(1)
    pub_reg = pub_reg.sort_values("total_pubs", ascending=True)
    fig = go.Figure()
    fig.add_trace(go.Bar(
        y=pub_reg["region"], x=pub_reg["total_pubs"], orientation="h",
        name="Total Publications", marker_color="#1f77b4",
    ))
    fig.update_layout(
        xaxis_title="Total Publications", template="plotly_white",
        height=380, margin=dict(t=10),
    )
    return _plotly_to_png(fig, width=1400, height=580)


# ────────────────────────────────────────────────────────────────────────────
# Network figures (matplotlib — same code path as dashboard)
# ────────────────────────────────────────────────────────────────────────────

def _fig_institution_network(df: pd.DataFrame) -> bytes:
    """Mirror app.py's `render_collab_matplotlib` with identical palette/layout."""
    G = nx.Graph()
    for _, row in df.iterrows():
        lead = str(row.get("institution", "")).strip().lower()
        if not lead or lead == "nan":
            continue
        orgs_raw = str(row.get("corrected_additional_orgs", "") or "")
        for org in orgs_raw.split(";"):
            org = org.strip().lower()
            if org and org != lead:
                if G.has_edge(lead, org):
                    G[lead][org]["weight"] += 1
                else:
                    G.add_edge(lead, org, weight=1)
    if G.number_of_edges() == 0:
        return _placeholder_png("No institution network data")
    degree_map = dict(G.degree())
    top_nodes = sorted(degree_map, key=lambda x: -degree_map[x])[:22]
    sub = G.subgraph(top_nodes).copy()
    sub.remove_edges_from([(u, v) for u, v, d in sub.edges(data=True) if d["weight"] < 3])
    sub.remove_nodes_from(list(nx.isolates(sub)))
    try:
        communities = list(nx.algorithms.community.greedy_modularity_communities(sub))
    except Exception:
        communities = [set(sub.nodes())]
    cmap = {n: i for i, c in enumerate(communities) for n in c}
    PALETTE = ["#2166ac", "#d6604d", "#4dac26", "#8073ac", "#f1a340", "#01665e"]

    n_comm = max(len(communities), 1)
    CLUSTER_R, INNER_R_BASE = 1.5, 0.42
    comm_centers = {}
    for i, cid in enumerate(sorted({cmap[n] for n in sub.nodes})):
        angle = -np.pi / 2 + 2 * np.pi * i / n_comm
        comm_centers[cid] = (CLUSTER_R * np.cos(angle), CLUSTER_R * np.sin(angle))
    comm_members = {}
    for n in sub.nodes:
        comm_members.setdefault(cmap.get(n, 0), []).append(n)
    pos = {}
    for cid, members in comm_members.items():
        cx_c, cy_c = comm_centers.get(cid, (0, 0))
        ordered = sorted(members, key=lambda n: -degree_map.get(n, 0))
        n = len(ordered)
        inner_r = INNER_R_BASE + 0.05 * n
        pos[ordered[0]] = (cx_c, cy_c)
        for j, node in enumerate(ordered[1:]):
            a = 2 * np.pi * j / (n - 1) if n > 2 else np.pi * j
            pos[node] = (cx_c + inner_r * np.cos(a), cy_c + inner_r * np.sin(a))
    pa = {n: list(p) for n, p in pos.items()}
    nl = list(pa)
    for _ in range(80):
        for i in range(len(nl)):
            for j in range(i + 1, len(nl)):
                a, b = nl[i], nl[j]
                dx = pa[b][0] - pa[a][0]; dy = pa[b][1] - pa[a][1]
                d = max((dx ** 2 + dy ** 2) ** 0.5, 1e-6)
                if d < 0.22:
                    push = (0.22 - d) / 2
                    ux, uy = dx / d, dy / d
                    pa[a][0] -= push * ux; pa[a][1] -= push * uy
                    pa[b][0] += push * ux; pa[b][1] += push * uy
    pos = {n: tuple(p) for n, p in pa.items()}

    max_w = max((d["weight"] for _, _, d in sub.edges(data=True)), default=1)
    max_deg = max(degree_map.values(), default=1)
    fig, ax = plt.subplots(figsize=(14, 11))
    fig.patch.set_facecolor("#ffffff")
    ax.set_facecolor("#f7f9f8"); ax.set_aspect("equal"); ax.axis("off")

    for u, v, d in sorted(sub.edges(data=True), key=lambda e: e[2]["weight"]):
        x0, y0 = pos[u]; x1, y1 = pos[v]
        cu, cv = cmap.get(u, 0), cmap.get(v, 0)
        w = d["weight"]
        if cu == cv:
            color = PALETTE[cu % len(PALETTE)]
            alpha = 0.18 + 0.32 * (w / max_w)
        else:
            color = "#999999"; alpha = 0.10 + 0.12 * (w / max_w)
        ax.plot([x0, x1], [y0, y1], "-", color=color,
                lw=0.5 + 3.5 * (w / max_w), alpha=alpha, zorder=1, solid_capstyle="round")
    for cid, (cx_c, cy_c) in comm_centers.items():
        ax.add_patch(plt.Circle((cx_c, cy_c), INNER_R_BASE + 0.28,
                                color=PALETTE[cid % len(PALETTE)],
                                alpha=0.06, zorder=0, linewidth=0))
    for n in sub.nodes:
        x, y = pos[n]
        ax.scatter([x], [y], s=180 + 820 * (degree_map.get(n, 1) / max_deg),
                   c=PALETTE[cmap.get(n, 0) % len(PALETTE)],
                   zorder=3, edgecolors="white", linewidths=2.8, alpha=0.94)

    def short_name(raw):
        ABBREV = {
            "university college london": "UCL",
            "king's college london": "King's College London",
            "imperial college london": "Imperial College London",
        }
        low = raw.lower()
        if low in ABBREV:
            return ABBREV[low]
        for prefix in ("the university of ", "university of ", "university college "):
            if low.startswith(prefix):
                return raw[len(prefix):].title()
        return raw.title().replace("University", "Univ.")

    for n in sub.nodes:
        x, y = pos[n]
        cid = cmap.get(n, 0)
        cx_c, cy_c = comm_centers.get(cid, (0, 0))
        dx, dy = x - cx_c, y - cy_c
        d = max((dx ** 2 + dy ** 2) ** 0.5, 1e-4)
        offset = 0.30 + 0.10 * (degree_map.get(n, 1) / max_deg)
        lx, ly = x + offset * dx / d, y + offset * dy / d
        ax.text(lx, ly, short_name(n), fontsize=14,
                fontweight="bold" if degree_map.get(n, 0) >= sorted(degree_map.values(), reverse=True)[min(4, len(degree_map) - 1)] else "normal",
                ha="center", va="center", zorder=5,
                bbox=dict(boxstyle="round,pad=0.28", facecolor="white",
                          edgecolor="#c0c0c0", alpha=0.90, linewidth=0.6))
    fig.tight_layout(pad=0.4)
    return _mpl_to_png(fig, dpi=220)


def _fig_topic_network(df: pd.DataFrame, ti: pd.DataFrame, *, top_n: int = 35, min_w: int = 4) -> bytes:
    """Mirror app.py topic network defaults (sus_only=False, kind=Topics+Subjects).

    Placeholder tags (e.g. 'See subject area') are already filtered at
    build_topic_index time. Funding per tag uses fractional attribution
    (project funding divided by its tag count) to remove bibliometric
    inflation from heavy-tagged multi-discipline grants.
    """
    ti = ti.assign(project_ref=ti["project_ref"].astype(str))
    df_ref = df.assign(project_ref=df["project_ref"].astype(str))
    tag_counts = ti.groupby("tag")["project_ref"].nunique().sort_values(ascending=False)
    top_tags = set(tag_counts.head(top_n).index)
    ti_f = ti[ti["tag"].isin(top_tags)]
    tags_per_project = ti.groupby("project_ref").size().to_dict()
    fund_lookup = df_ref.set_index("project_ref")["fund_value"].to_dict()
    tag_funding: dict[str, float] = {}
    for tag, ref in (
        ti_f.drop_duplicates(["tag", "project_ref"])[["tag", "project_ref"]].itertuples(index=False)
    ):
        n_tags = max(tags_per_project.get(ref, 1), 1)
        proj_fund = float(fund_lookup.get(ref, 0) or 0)
        tag_funding[tag] = tag_funding.get(tag, 0.0) + proj_fund / n_tags
    G = nx.Graph()
    for t in top_tags:
        G.add_node(t, count=int(tag_counts[t]), funding=tag_funding.get(t, 0.0))
    for tags in ti_f.groupby("project_ref")["tag"].apply(lambda s: list(set(s))):
        if len(tags) < 2:
            continue
        for i in range(len(tags)):
            for j in range(i + 1, len(tags)):
                u, v = tags[i], tags[j]
                if G.has_edge(u, v):
                    G[u][v]["weight"] += 1
                else:
                    G.add_edge(u, v, weight=1)
    G.remove_edges_from([(u, v) for u, v, d in G.edges(data=True) if d["weight"] < min_w])
    G.remove_nodes_from(list(nx.isolates(G)))
    if G.number_of_nodes() == 0:
        return _placeholder_png("No topic co-occurrences")
    try:
        communities = list(nx.algorithms.community.greedy_modularity_communities(G))
    except Exception:
        communities = [set(G.nodes)]
    cmap = {n: i for i, c in enumerate(communities) for n in c}
    PALETTE = ["#2166ac", "#d6604d", "#4dac26", "#8073ac", "#f1a340",
               "#01665e", "#c51b7d", "#5aae61", "#8c510a", "#7570b3"]
    pos = nx.spring_layout(G, seed=42, k=1.4 / (max(len(G), 1) ** 0.5), iterations=140, weight="weight")
    counts = {n: G.nodes[n].get("count", 1) for n in G.nodes}
    max_c = max(counts.values(), default=1)
    fundings = {n: G.nodes[n].get("funding", 0.0) for n in G.nodes}
    max_f = max(fundings.values(), default=1) or 1
    max_w = max((d["weight"] for _, _, d in G.edges(data=True)), default=1)
    fig, ax = plt.subplots(figsize=(15, 11))
    fig.patch.set_facecolor("#ffffff"); ax.set_facecolor("#f7f9f8")
    ax.set_aspect("equal"); ax.axis("off")
    for u, v, d in sorted(G.edges(data=True), key=lambda e: e[2]["weight"]):
        x0, y0 = pos[u]; x1, y1 = pos[v]
        cu, cv = cmap.get(u, 0), cmap.get(v, 0)
        w = d["weight"]
        if cu == cv:
            color = PALETTE[cu % len(PALETTE)]; alpha = 0.40 + 0.45 * (w / max_w)
        else:
            color = "#888888"; alpha = 0.25 + 0.35 * (w / max_w)
        lw = 2.0 + 11.0 * (w / max_w) ** 0.6
        ax.plot([x0, x1], [y0, y1], "-", color=color,
                lw=lw, alpha=alpha, zorder=1, solid_capstyle="round")
    for n in G.nodes:
        x, y = pos[n]
        f_norm = (fundings[n] / max_f) ** 0.5 if max_f else 0
        ax.scatter([x], [y], s=220 + 2400 * f_norm,
                   c=PALETTE[cmap.get(n, 0) % len(PALETTE)],
                   edgecolors="white", linewidths=2.4, zorder=3)
    for n in G.nodes:
        x, y = pos[n]
        short = n if len(n) <= 28 else n[:26] + "…"
        f_m = fundings[n] / 1e6
        if f_m >= 100:
            f_str = f"£{f_m:.0f}M"
        elif f_m >= 1:
            f_str = f"£{f_m:.1f}M"
        else:
            f_str = f"£{f_m * 1000:.0f}k"
        label = f"{short}\n{f_str} · {counts[n]} proj"
        ax.text(x, y, label, fontsize=9, ha="center", va="center",
                fontweight="bold", zorder=5,
                bbox=dict(boxstyle="round,pad=0.24", facecolor="white",
                          edgecolor="#c0c0c0", alpha=0.92, linewidth=0.6))
    fig.tight_layout(pad=0.4)
    return _mpl_to_png(fig, dpi=220)


def _fig_pi_network(pp: pd.DataFrame, *, top_n: int = 40, min_w: int = 1) -> bytes:
    pp = pp[pp["role"].isin({"PRINCIPAL_INVESTIGATOR", "CO_INVESTIGATOR"})].copy()
    pp["project_ref"] = pp["project_ref"].astype(str)
    if pp.empty:
        return _placeholder_png("No PI data")
    G_full = nx.Graph()
    for ppl in pp.groupby("project_ref")["full_name"].apply(lambda s: list(set(s))):
        if len(ppl) < 2:
            continue
        for i in range(len(ppl)):
            for j in range(i + 1, len(ppl)):
                u, v = ppl[i], ppl[j]
                if G_full.has_edge(u, v):
                    G_full[u][v]["weight"] += 1
                else:
                    G_full.add_edge(u, v, weight=1)
    G_full.remove_edges_from([(u, v) for u, v, d in G_full.edges(data=True) if d["weight"] < min_w])
    if G_full.number_of_edges() == 0:
        return _placeholder_png("No PI edges")
    deg = dict(G_full.degree(weight="weight"))
    top = [n for n, _ in sorted(deg.items(), key=lambda kv: -kv[1])[:top_n]]
    sub = G_full.subgraph(top).copy()
    sub.remove_nodes_from(list(nx.isolates(sub)))
    if sub.number_of_nodes() == 0:
        return _placeholder_png("No PI subgraph")
    try:
        communities = list(nx.algorithms.community.greedy_modularity_communities(sub))
    except Exception:
        communities = [set(sub.nodes)]
    cmap = {n: i for i, c in enumerate(communities) for n in c}
    PALETTE = ["#2166ac", "#d6604d", "#4dac26", "#8073ac", "#f1a340",
               "#01665e", "#c51b7d", "#5aae61", "#8c510a", "#7570b3"]
    n_nodes = max(len(sub), 1)
    pos = nx.spring_layout(sub, seed=7, k=2.6 / (n_nodes ** 0.5), iterations=240, weight="weight")
    min_sep = max(0.18, 1.3 / (n_nodes ** 0.5))
    pa = {n: list(p) for n, p in pos.items()}
    nl = list(pa)
    for _ in range(120):
        moved = False
        for i in range(len(nl)):
            for j in range(i + 1, len(nl)):
                a, b = nl[i], nl[j]
                dx = pa[b][0] - pa[a][0]; dy = pa[b][1] - pa[a][1]
                d = max((dx ** 2 + dy ** 2) ** 0.5, 1e-6)
                if d < min_sep:
                    push = (min_sep - d) / 2
                    ux, uy = dx / d, dy / d
                    pa[a][0] -= push * ux; pa[a][1] -= push * uy
                    pa[b][0] += push * ux; pa[b][1] += push * uy
                    moved = True
        if not moved:
            break
    pos = {n: tuple(p) for n, p in pa.items()}
    counts = {n: int(pp.loc[pp["full_name"] == n, "project_ref"].nunique()) for n in sub.nodes}
    max_c = max(counts.values(), default=1)
    max_w = max((d["weight"] for _, _, d in sub.edges(data=True)), default=1)
    fig, ax = plt.subplots(figsize=(17, 13))
    fig.patch.set_facecolor("#ffffff"); ax.set_facecolor("#f7f9f8")
    ax.set_aspect("equal"); ax.axis("off")
    comm_members = {}
    for n, cid in cmap.items():
        comm_members.setdefault(cid, []).append(n)
    comm_centers = {cid: (sum(pos[m][0] for m in ms) / len(ms),
                           sum(pos[m][1] for m in ms) / len(ms))
                    for cid, ms in comm_members.items()}
    for u, v, d in sorted(sub.edges(data=True), key=lambda e: e[2]["weight"]):
        x0, y0 = pos[u]; x1, y1 = pos[v]
        cu, cv = cmap.get(u, 0), cmap.get(v, 0)
        w = d["weight"]
        if cu == cv:
            color = PALETTE[cu % len(PALETTE)]; alpha = 0.55 + 0.35 * (w / max_w)
        else:
            color = "#7a7a7a"; alpha = 0.40 + 0.30 * (w / max_w)
        ax.plot([x0, x1], [y0, y1], "-", color=color,
                lw=1.0 + 3.0 * (w / max_w), alpha=alpha, zorder=1, solid_capstyle="round")
    for n in sub.nodes:
        x, y = pos[n]
        ax.scatter([x], [y], s=180 + 620 * (counts[n] / max_c),
                   c=PALETTE[cmap.get(n, 0) % len(PALETTE)],
                   edgecolors="white", linewidths=2.2, zorder=3)
    sorted_counts = sorted(counts.values(), reverse=True)
    top10 = sorted_counts[min(9, len(sorted_counts) - 1)] if sorted_counts else 1
    for n in sub.nodes:
        x, y = pos[n]
        cid = cmap.get(n, 0); cx_c, cy_c = comm_centers.get(cid, (0, 0))
        dx, dy = x - cx_c, y - cy_c
        dist = (dx ** 2 + dy ** 2) ** 0.5
        offset = max(0.06, 0.13 - 0.04)
        if dist < 1e-3:
            lx, ly = x, y - offset
        else:
            lx, ly = x + offset * dx / dist, y + offset * dy / dist
        label = n if len(n) <= 24 else n[:22] + "…"
        bold = counts[n] >= top10
        ax.text(lx, ly, label, fontsize=9 if bold else 8,
                fontweight="bold" if bold else "normal",
                ha="center", va="center", zorder=5,
                bbox=dict(boxstyle="round,pad=0.18", facecolor="white",
                          edgecolor="#c8c8c8", alpha=0.93, linewidth=0.5))
    fig.tight_layout(pad=0.4)
    return _mpl_to_png(fig, dpi=220)


def _fig_org_project_bipartite(df_all: pd.DataFrame, op: pd.DataFrame,
                                *, top_orgs_n: int = 14, top_projects_n: int = 30) -> bytes:
    op = op.assign(project_ref=op["project_ref"].astype(str))
    df_all = df_all.assign(project_ref=df_all["project_ref"].astype(str))
    sus_refs = set(df_all.loc[df_all["is_sustainability"], "project_ref"])
    op_s = op[op["project_ref"].isin(sus_refs)].copy()
    per_proj = op_s.groupby("project_ref")["org_name"].nunique()
    multi = per_proj[per_proj >= 2].index
    op_s = op_s[op_s["project_ref"].isin(multi)]
    if op_s.empty:
        return _placeholder_png("No multi-org sustainability projects")
    ranked = op_s.groupby("org_name")["project_ref"].nunique().sort_values(ascending=False)
    top_orgs = ranked.head(top_orgs_n).index.tolist()
    op_focus = op_s[op_s["org_name"].isin(top_orgs)]
    proj_meta = (
        df_all.loc[df_all["project_ref"].isin(op_focus["project_ref"].unique()),
                   ["project_ref", "title", "funder", "fund_value"]]
              .drop_duplicates("project_ref").set_index("project_ref")
    )
    proj_score = op_focus.groupby("project_ref")["org_name"].nunique().rename("top_org_count")
    proj_rank = (proj_score.to_frame().join(proj_meta[["fund_value"]], how="left")
                 .fillna({"fund_value": 0})
                 .sort_values(["top_org_count", "fund_value"], ascending=[False, False]))
    top_proj_refs = proj_rank.head(top_projects_n).index.tolist()
    op_final = op_focus[op_focus["project_ref"].isin(top_proj_refs)]
    G = nx.Graph()
    for org in top_orgs:
        G.add_node(("org", org), kind="org", label=org)
    for ref in top_proj_refs:
        if ref in proj_meta.index:
            meta = proj_meta.loc[ref]
            title = str(meta["title"]) if pd.notna(meta["title"]) else ref
            funder = str(meta["funder"]) if pd.notna(meta["funder"]) else "Unknown"
        else:
            title, funder = ref, "Unknown"
        G.add_node(("proj", ref), kind="proj", label=title, funder=funder, ref=ref)
    for _, row in op_final.iterrows():
        G.add_edge(("org", row["org_name"]), ("proj", row["project_ref"]))
    for org in list(top_orgs):
        if ("org", org) in G and G.degree(("org", org)) == 0:
            G.remove_node(("org", org))

    org_nodes = [n for n in G.nodes if G.nodes[n].get("kind") == "org"]
    proj_nodes = [n for n in G.nodes if G.nodes[n].get("kind") == "proj"]
    if not org_nodes or not proj_nodes:
        return _placeholder_png("No bipartite edges")

    FUNDER_PALETTE = ["#1f77b4", "#d62728", "#2ca02c", "#ff7f0e", "#9467bd",
                      "#8c564b", "#e377c2", "#17becf", "#bcbd22", "#7f7f7f"]
    funders = sorted({G.nodes[n].get("funder", "Unknown") for n in proj_nodes})
    funder_color = {f: FUNDER_PALETTE[i % len(FUNDER_PALETTE)] for i, f in enumerate(funders)}

    org_sorted = sorted(org_nodes, key=lambda n: -G.degree(n))
    proj_sorted = sorted(proj_nodes, key=lambda n: (G.nodes[n].get("funder", "zzz"), -G.degree(n)))

    def col_pos(nodes, x):
        m = len(nodes)
        if m == 1:
            return {nodes[0]: (x, 0)}
        return {n: (x, 1.0 - 2.0 * i / (m - 1)) for i, n in enumerate(nodes)}

    pos = {}
    pos.update(col_pos(org_sorted, -1.0))
    pos.update(col_pos(proj_sorted, 1.0))
    n_rows = max(len(org_sorted), len(proj_sorted))
    fig_h = max(8.0, 0.36 * n_rows)
    fig, ax = plt.subplots(figsize=(16, fig_h))
    fig.patch.set_facecolor("#ffffff"); ax.set_facecolor("#f7f9f8"); ax.axis("off")
    for u, v in G.edges():
        proj_node = u if G.nodes[u].get("kind") == "proj" else v
        c = funder_color.get(G.nodes[proj_node].get("funder", "Unknown"), "#888")
        x0, y0 = pos[u]; x1, y1 = pos[v]
        ax.plot([x0, x1], [y0, y1], "-", color=c, alpha=0.22, lw=0.9, zorder=1)
    max_deg = max((G.degree(n) for n in org_sorted), default=1)
    for n in org_sorted:
        x, y = pos[n]
        size = 240 + 700 * (G.degree(n) / max_deg)
        ax.scatter([x], [y], s=size, c="#0b3d91", edgecolors="white", linewidths=2.2, zorder=3)
        label = G.nodes[n]["label"].title()
        if len(label) > 38:
            label = label[:36] + "…"
        ax.text(x - 0.08, y, label, ha="right", va="center",
                fontsize=11, fontweight="bold", zorder=5)
    for n in proj_sorted:
        x, y = pos[n]
        c = funder_color.get(G.nodes[n].get("funder", "Unknown"), "#888")
        ax.scatter([x], [y], s=110, c=c, edgecolors="white", linewidths=1.2, zorder=3, marker="s")
        title = G.nodes[n].get("label", "")
        short = title if len(title) <= 60 else title[:57] + "…"
        ax.text(x + 0.06, y, short, ha="left", va="center", fontsize=9, zorder=5)
    ax.text(-1.0, 1.07, "Organisations", ha="center", va="bottom",
            fontsize=15, fontweight="bold", color="#0b3d91")
    ax.text(1.0, 1.07, "Sustainability Projects", ha="center", va="bottom",
            fontsize=15, fontweight="bold", color="#333")
    ax.set_xlim(-2.6, 2.6); ax.set_ylim(-1.18, 1.18)
    fig.tight_layout(pad=0.4)
    return _mpl_to_png(fig, dpi=220)


def _fig_theme_chord(df: pd.DataFrame) -> bytes:
    """Mirror app.py render_theme_network — chord-layout with edge weight labels."""
    from matplotlib.patches import FancyArrowPatch
    import math
    G = nx.Graph()
    for theme in SUSTAINABILITY_THEMES:
        G.add_node(theme)
    for s in df["sustainability_themes"].fillna(""):
        themes = [t.strip() for t in str(s).split(";") if t.strip()]
        for i in range(len(themes)):
            for j in range(i + 1, len(themes)):
                u, v = themes[i], themes[j]
                if u in G and v in G:
                    if G.has_edge(u, v):
                        G[u][v]["weight"] += 1
                    else:
                        G.add_edge(u, v, weight=1)
    for theme in SUSTAINABILITY_THEMES:
        col = f"theme_{theme.lower().replace(' ', '_').replace('&', 'and')}"
        if col in df.columns:
            G.nodes[theme]["count"] = int(df[col].sum())
    if G.number_of_edges() == 0:
        return _placeholder_png("No theme co-occurrences")
    PALETTE = ["#1b7837", "#5aae61", "#762a83", "#9970ab",
               "#d6604d", "#f1a340", "#4393c3", "#2166ac"]
    try:
        communities = list(nx.algorithms.community.greedy_modularity_communities(G))
    except Exception:
        communities = [set(G.nodes())]
    cmap = {n: i for i, c in enumerate(communities) for n in c}
    comm_groups = {}
    for n, cid in cmap.items():
        comm_groups.setdefault(cid, []).append(n)
    node_order = []
    for cid in sorted(comm_groups):
        node_order.extend(sorted(comm_groups[cid], key=lambda n: -G.nodes[n].get("count", 0)))
    n_nodes = len(node_order)
    R = 1.0
    pos = {}
    for i, n in enumerate(node_order):
        theta = math.pi / 2 - 2 * math.pi * i / n_nodes
        pos[n] = (R * math.cos(theta), R * math.sin(theta))
    max_w = max((d["weight"] for _, _, d in G.edges(data=True)), default=1)
    counts = {n: G.nodes[n].get("count", 1) for n in G.nodes}
    max_c = max(counts.values(), default=1)
    fig, ax = plt.subplots(figsize=(14, 11))
    fig.patch.set_facecolor("#ffffff"); ax.set_facecolor("#fbfcfa")
    ax.set_aspect("equal"); ax.axis("off")
    edges_sorted = sorted(G.edges(data=True), key=lambda e: e[2]["weight"])
    for u, v, d in edges_sorted:
        cu, cv = cmap.get(u, 0), cmap.get(v, 0)
        w = d["weight"]
        if cu == cv:
            color = PALETTE[cu % len(PALETTE)]; alpha = 0.55 + 0.35 * (w / max_w)
        else:
            color = "#888888"; alpha = 0.30 + 0.40 * (w / max_w)
        ax.add_patch(FancyArrowPatch(
            pos[u], pos[v], connectionstyle="arc3,rad=0.18",
            arrowstyle="-", linewidth=1.0 + 8.5 * (w / max_w),
            color=color, alpha=alpha, zorder=1, capstyle="round",
        ))
    median_w = sorted([e[2]["weight"] for e in edges_sorted])[len(edges_sorted) // 2]
    for u, v, d in edges_sorted:
        w = d["weight"]
        x0, y0 = pos[u]; x1, y1 = pos[v]
        mx, my = (x0 + x1) / 2, (y0 + y1) / 2
        dx, dy = x1 - x0, y1 - y0
        length = (dx ** 2 + dy ** 2) ** 0.5 or 1e-6
        nx_v, ny_v = -dy / length, dx / length
        curve_offset = 0.18 * length * 0.5
        lx, ly = mx + nx_v * curve_offset, my + ny_v * curve_offset
        strong = w >= median_w
        ax.text(lx, ly, f"{w}",
                fontsize=10 if strong else 8,
                fontweight="bold" if strong else "normal",
                ha="center", va="center",
                color="#222" if strong else "#555", zorder=4,
                bbox=dict(boxstyle="round,pad=0.18", facecolor="#ffffff",
                          edgecolor="#bbb" if strong else "#dcdcdc",
                          alpha=0.93 if strong else 0.80, linewidth=0.5))
    for n in G.nodes:
        x, y = pos[n]
        ax.scatter([x], [y], s=800 + 2200 * (counts[n] / max_c),
                   c=PALETTE[cmap.get(n, 0) % len(PALETTE)],
                   edgecolors="white", linewidths=3.4, zorder=3, alpha=0.95)
    for n in G.nodes:
        x, y = pos[n]
        theta = math.atan2(y, x)
        label_r = 1.32
        lx, ly = label_r * math.cos(theta), label_r * math.sin(theta)
        if math.cos(theta) > 0:
            ha = "left"; rot = math.degrees(theta)
        else:
            ha = "right"; rot = math.degrees(theta) + 180
        if abs(math.cos(theta)) < 0.15:
            rot = 0; ha = "center"
        ax.text(lx, ly, f"{n}\n({counts[n]} projects)", fontsize=11, ha=ha, va="center",
                fontweight="bold", color="#1a3d2b",
                rotation=rot, rotation_mode="anchor", zorder=5,
                bbox=dict(boxstyle="round,pad=0.32", facecolor="white",
                          edgecolor="#b8b8b8", alpha=0.94, linewidth=0.7))
    ax.set_xlim(-2.3, 2.3); ax.set_ylim(-1.7, 1.7)
    fig.tight_layout(pad=0.4)
    return _mpl_to_png(fig, dpi=240)


def _fig_theme_matrix(df: pd.DataFrame) -> bytes:
    themes = list(SUSTAINABILITY_THEMES.keys())
    theme_cols = {t: f"theme_{t.lower().replace(' ', '_').replace('&', 'and')}" for t in themes}
    available = [t for t in themes if theme_cols[t] in df.columns]
    if not available:
        return _placeholder_png("No theme columns")
    M = np.zeros((len(available), len(available)), dtype=int)
    for i, a in enumerate(available):
        col_a = df[theme_cols[a]]
        for j, b in enumerate(available):
            col_b = df[theme_cols[b]]
            if i == j:
                M[i, j] = int(col_a.sum())
            else:
                M[i, j] = int((col_a & col_b).sum())
    fig, ax = plt.subplots(figsize=(11, 8.5))
    fig.patch.set_facecolor("#ffffff")
    im = ax.imshow(M, cmap=plt.cm.YlGn, aspect="equal")
    for i in range(len(available)):
        ax.add_patch(plt.Rectangle((i - 0.5, i - 0.5), 1, 1, fill=False,
                                   edgecolor="#333", linewidth=1.4, zorder=4))
    vmax = M.max() if M.max() > 0 else 1
    for i in range(len(available)):
        for j in range(len(available)):
            v = M[i, j]
            if v == 0:
                continue
            color = "white" if v / vmax > 0.55 else "#1a3d2b"
            ax.text(j, i, str(v), ha="center", va="center",
                    fontsize=10, fontweight="bold" if i == j else "normal", color=color)
    ax.set_xticks(range(len(available))); ax.set_yticks(range(len(available)))
    ax.set_xticklabels(available, rotation=40, ha="right", fontsize=10)
    ax.set_yticklabels(available, fontsize=10)
    ax.set_title("Theme co-occurrence matrix  ·  diagonal = total projects per theme",
                 fontsize=12, color="#333", pad=12)
    fig.colorbar(im, ax=ax, fraction=0.04, pad=0.02).set_label("Shared projects", fontsize=10)
    fig.tight_layout(pad=0.4)
    return _mpl_to_png(fig, dpi=220)


def _fig_funder_inst_bipartite(df: pd.DataFrame) -> bytes:
    top_insts = (df.groupby("institution").size()
                 .sort_values(ascending=False).head(20).index.tolist())
    top_funders = df["funder"].dropna().value_counts().index.tolist()
    edge_df = (df[df["institution"].isin(top_insts)]
               .groupby(["funder", "institution"]).size().reset_index(name="count"))
    if edge_df.empty:
        return _placeholder_png("No funder-institution edges")
    funder_counts = df["funder"].value_counts().to_dict()
    inst_counts = df[df["institution"].isin(top_insts)].groupby("institution").size().to_dict()
    max_w = edge_df["count"].max() or 1
    max_fc = max(funder_counts.values(), default=1)
    max_ic = max(inst_counts.values(), default=1)
    fy = {f: 1.0 - i / max(len(top_funders) - 1, 1) for i, f in enumerate(top_funders)}
    iy = {inst: 1.0 - i / max(len(top_insts) - 1, 1) for i, inst in enumerate(top_insts)}
    X_F, X_I = 0.0, 1.0
    fig, ax = plt.subplots(figsize=(18, 14))
    fig.patch.set_facecolor("#ffffff"); ax.set_facecolor("#f7f9f8")
    ax.set_xlim(-0.62, 1.62); ax.set_ylim(-0.06, 1.10); ax.axis("off")
    for _, row in edge_df.sort_values("count").iterrows():
        f, inst, w = row["funder"], row["institution"], row["count"]
        if f not in fy or inst not in iy:
            continue
        color = COUNCIL_COLORS.get(f, "#aaa")
        ax.plot([X_F, X_I], [fy[f], iy[inst]], "-", color=color,
                lw=0.5 + 5.5 * (w / max_w), alpha=0.14 + 0.48 * (w / max_w),
                zorder=1, solid_capstyle="round")
    for f in top_funders:
        ax.scatter([X_F], [fy[f]], s=80 + 320 * (funder_counts.get(f, 1) / max_fc),
                   c=COUNCIL_COLORS.get(f, "#aaa"), edgecolors="white",
                   linewidths=2, zorder=3, alpha=0.93)
        ax.text(X_F - 0.04, fy[f], f, ha="right", va="center", fontsize=18, color="#1a2a3a")
    for inst in top_insts:
        ax.scatter([X_I], [iy[inst]], s=80 + 320 * (inst_counts.get(inst, 1) / max_ic),
                   c="#2ca02c", edgecolors="white", linewidths=2, zorder=3, alpha=0.93)
        ax.text(X_I + 0.04, iy[inst], inst.title()[:40], ha="left", va="center",
                fontsize=18, color="#1a5438")
    ax.text(X_F, 1.07, "Funding Councils", ha="center", va="bottom",
            fontsize=22, fontweight="bold", color="#1f4e79")
    ax.text(X_I, 1.07, "Lead Institutions (Top 20)", ha="center", va="bottom",
            fontsize=22, fontweight="bold", color="#1a5438")
    fig.tight_layout(pad=0.3)
    return _mpl_to_png(fig, dpi=300)


# ────────────────────────────────────────────────────────────────────────────
# Top-level build
# ────────────────────────────────────────────────────────────────────────────

def build_pptx_bytes(df: pd.DataFrame) -> bytes:
    op = get_or_build_org_index().assign(project_ref=lambda d: d["project_ref"].astype(str))
    ti = get_or_build_topic_index().assign(project_ref=lambda d: d["project_ref"].astype(str))
    pp = get_or_build_person_index().assign(project_ref=lambda d: d["project_ref"].astype(str))
    df = df.copy()
    df["project_ref"] = df["project_ref"].astype(str)

    prs = _new_presentation()

    # ── Title + headline metrics ───────────────────────────────────────────
    _add_title_slide(prs,
        title="UKRI Sustainability Research Funding",
        subtitle=f"A dashboard for Nature Sustainability — generated {date.today().isoformat()}\n"
                 f"Source: UKRI Gateway-to-Research (JSON) + curated Excel registry.",
        notes=("This deck mirrors every chart, table, and network in the live dashboard. "
               "Plotly figures are exported through kaleido so the colours, scales, and "
               "aspect ratios match the web app exactly."),
    )

    total_projects = len(df)
    total_funding = df["fund_value"].sum()
    sus_projects = int(df["is_sustainability"].sum())
    sus_share = sus_projects / total_projects if total_projects else 0
    sus_funding = df.loc[df["is_sustainability"], "fund_value"].sum()
    yr_min = int(df["start_year"].min()) if df["start_year"].notna().any() else 0
    yr_max = int(df["start_year"].max()) if df["start_year"].notna().any() else 0
    _add_metrics_slide(prs, "Headline metrics",
        [
            ("Projects analysed", f"{total_projects:,}"),
            ("Total UKRI funding", f"£{total_funding/1e9:.2f}B"),
            ("Sustainability-classified projects", f"{sus_projects:,}  ({sus_share*100:.1f}%)"),
            ("Sustainability funding", f"£{sus_funding/1e9:.2f}B"),
            ("Time window", f"{yr_min} – {yr_max}"),
            ("Funding councils represented", f"{df['funder'].nunique()}"),
            ("Lead institutions", f"{df['institution'].nunique():,}"),
            ("JSON projects indexed", f"{ti['project_ref'].nunique():,}"),
        ],
        notes=("Headline numbers match the four metric cards at the top of the Overview "
               "tab plus a few derived counts. Use these to open the talk."),
    )

    # ── Tab 1 — Overview ──────────────────────────────────────────────────
    _add_section_divider(prs, "1.  Overview", notes="Tab 1 — funding totals over time, by council, by status.")
    _add_image_slide(prs, "Annual UKRI funding (£B) by research council",
                     _fig_annual_funding(df),
                     notes="Stacked bar of annual funding by start year, broken out per council. Colours match the dashboard's council palette.")
    _add_image_slide(prs, "Funding share by council",
                     _fig_funding_share_donut(df),
                     notes="Donut chart of total funding share per council across the entire portfolio.")
    _add_image_slide(prs, "Sustainability funding trend",
                     _fig_sustainability_trend(df),
                     notes="Stacked area: sustainability (green) sits on top of non-sustainability (light blue). Shows the rising sustainability share.")
    _add_image_slide(prs, "Project count by year & status",
                     _fig_project_count_by_year(df),
                     notes="Annual project starts coloured by current status — Active vs Closed.")
    _add_image_slide(prs, "Grant size distribution by council",
                     _fig_grant_size_box(df),
                     notes="Per-council box plot of grant value on a log £M axis. Outliers hidden for readability.")

    # ── Tab 2 — Geographic distribution ────────────────────────────────────
    _add_section_divider(prs, "2.  Geographic distribution",
                        notes="Tab 2 — funding and projects by region and nation.")
    _add_image_slide(prs, "Total funding by UK region (£M)",
                     _fig_region_funding(df),
                     notes="Regional ranking by total funding with £M annotations.")
    _add_image_slide(prs, "Projects by region",
                     _fig_projects_by_region(df),
                     notes="Region ranking by project count, with bar colour encoding the %-sustainability share.")
    _add_image_slide(prs, "Funding by UK nation",
                     _fig_nation_funding(df),
                     notes="Stacked bar per nation — sustainability funding green, all other funding light blue.")
    _add_image_slide(prs, "Regional funding over time (£M) — heatmap",
                     _fig_regional_heatmap(df),
                     notes="Region × year heatmap of funding totals. Bright bands mark cohort years with portfolio-wide spend.")

    # ── Tab 3 — Institutions ───────────────────────────────────────────────
    _add_section_divider(prs, "3.  Institutions",
                        notes="Tab 3 — institution rankings, council ties, departments.")
    _add_image_slide(prs, "Top 20 institutions by total funding",
                     _fig_top_institutions(df),
                     notes="Top-20 by funding; bar colour encodes the institution's %-sustainability share.")
    _add_image_slide(prs, "Institution: funding vs. projects vs. publications",
                     _fig_institution_bubble(df),
                     notes="Bubble chart, top 50 institutions: x = project count, y = total funding, size = total publications, colour = % sustainability.")
    _add_image_slide(prs, "Department-level funding (top 15)",
                     _fig_department_funding(df),
                     notes="Top-15 departments by funding (excludes 'Unlisted' rows).")
    _add_image_slide(prs, "Institution × funding council matrix (£M)",
                     _fig_institution_council_matrix(df),
                     notes="£M flowing from each council to each top-15 institution. Annotated heatmap.")

    # ── Tab 4 — Sustainability themes ─────────────────────────────────────
    _add_section_divider(prs, "4.  Sustainability themes",
                        notes="Tab 4 — the eight themes, their volumes, evolutions, and intersections.")
    _add_image_slide(prs, "Projects per sustainability theme",
                     _fig_themes_bar(df),
                     notes="Bar of projects per theme, colour encoding £M funding per theme.")
    _add_image_slide(prs, "Sustainability funding by theme (£M)",
                     _fig_theme_funding_pie(df),
                     notes="Donut of £M sustainability funding by theme.")
    _add_image_slide(prs, "Sustainability theme evolution (projects per year)",
                     _fig_theme_evolution(df),
                     notes="Line per theme showing projects per start year.")
    _add_image_slide(prs, "Sustainability theme × funding council (project count)",
                     _fig_theme_council_heatmap(df),
                     notes="Annotated heatmap revealing which councils dominate each theme.")
    _add_image_slide(prs, "Sustainability theme × region (project count)",
                     _fig_theme_region_heatmap(df),
                     notes="Annotated heatmap showing the regional distribution of each theme.")
    _add_image_slide(prs, "Most frequent research topics in sustainability projects",
                     _fig_research_topics(df),
                     notes="Top-25 raw research-topic tags across sustainability-classified projects.")

    # ── Tab 5 — Collaboration networks ─────────────────────────────────────
    _add_section_divider(prs, "5.  Collaboration",
                        notes="Tab 5 — collaboration counts, distributions, and the network graphs.")
    _add_image_slide(prs, "Collaboration size distribution",
                     _fig_collab_histogram(df),
                     notes="Histogram of partner-organisation counts per project.")
    _add_image_slide(prs, "Collaboration size by funding council (box plot)",
                     _fig_collab_box_by_council(df),
                     notes="Per-council distribution of partner counts.")
    _add_image_slide(prs, "Average collaboration size by region & sustainability status",
                     _fig_collab_region_bar(df),
                     notes="Grouped bar per region: average partners for sustainability projects vs. everything else.")
    _add_image_slide(prs, "Collaboration size vs. grant value",
                     _fig_collab_vs_grant_scatter(df),
                     notes="Scatter coloured by council, with OLS trendline (red).")
    _add_image_slide(prs, "Top 20 most-collaborative lead institutions (avg. partners)",
                     _fig_top_collaborative_institutions(df),
                     notes="Average partner count per project across institutions with ≥3 projects.")
    _add_image_slide(prs, "Institution co-collaboration network (Top 22 hubs)",
                     _fig_institution_network(df),
                     notes="Matplotlib network from the dashboard. Communities by greedy modularity; node size ∝ partner count.")
    _add_image_slide(prs, "Organisation–Project network (sustainability collaborations)",
                     _fig_org_project_bipartite(df, op),
                     notes="Bipartite view of top organisations × multi-org sustainability projects (data from JSON `organisationRoles`).")
    _add_image_slide(prs, "Research topic co-occurrence network",
                     _fig_topic_network(df, ti),
                     notes="Topic-topic co-occurrence graph from `researchTopics` + `researchSubjects` JSON fields.")
    _add_image_slide(prs, "Principal Investigator collaboration network",
                     _fig_pi_network(pp),
                     notes="Investigators who co-PI/Co-I together; top-40 by weighted collaboration degree.")
    _add_image_slide(prs, "Sustainability theme co-occurrence network (chord)",
                     _fig_theme_chord(df),
                     notes="Chord-layout co-occurrence of the eight sustainability themes.")
    _add_image_slide(prs, "Sustainability theme co-occurrence matrix",
                     _fig_theme_matrix(df),
                     notes="Matrix view of the same data — diagonal shows total projects per theme.")

    # Top theme pairings table — built the same way as the dashboard
    pair_counts: dict[tuple[str, str], int] = {}
    for s in df["sustainability_themes"].fillna(""):
        themes = sorted({t.strip() for t in str(s).split(";") if t.strip()})
        for i in range(len(themes)):
            for j in range(i + 1, len(themes)):
                key = (themes[i], themes[j])
                pair_counts[key] = pair_counts.get(key, 0) + 1
    top_pairs = sorted(pair_counts.items(), key=lambda kv: -kv[1])[:8]
    if top_pairs:
        _add_table_slide(
            prs, "Top theme pairings (shared projects)",
            ["Theme A", "Theme B", "Shared projects"],
            [[a, b, str(n)] for (a, b), n in top_pairs],
            notes=(
                "Top eight cross-theme pairings, ranked by number of projects "
                "classified under both. Mirrors the table shown beneath the chord "
                "view of the theme co-occurrence network in the dashboard."
            ),
        )
    _add_image_slide(prs, "Funder–Institution bipartite network",
                     _fig_funder_inst_bipartite(df),
                     notes="Bipartite chart linking funding councils to the top-20 lead institutions by shared projects.")

    # ── Tab 6 — Research impact ────────────────────────────────────────────
    _add_section_divider(prs, "6.  Research impact",
                        notes="Tab 6 — publication output across the portfolio.")
    _add_image_slide(prs, "Publications per project by council",
                     _fig_pubs_per_project_by_council(df),
                     notes="Average publications per project per council.")
    _add_image_slide(prs, "Publications: sustainability vs other projects",
                     _fig_pubs_histogram(df),
                     notes="Overlaid histograms of publication counts, log y-axis.")

    if "publication_count" in df.columns:
        top_pub = (df[df["publication_count"] > 0]
                   .sort_values("publication_count", ascending=False)
                   .head(20)[["title", "funder", "institution",
                              "start_year", "publication_count", "fund_value"]])
        _add_table_slide(prs, "Top 20 projects by publication output",
            ["Title", "Council", "Institution", "Year", "Publications", "Grant (£M)"],
            [[(str(r["title"])[:55] + ("…" if len(str(r["title"])) > 55 else "")),
              str(r["funder"])[:24],
              str(r["institution"])[:32],
              f"{int(r['start_year']) if pd.notna(r['start_year']) else '—'}",
              f"{int(r['publication_count'])}",
              f"{(r['fund_value'] or 0)/1e6:.1f}"] for _, r in top_pub.iterrows()],
            notes="Top-20 projects by total publications recorded in JSON output sections.",
        )

    _add_image_slide(prs, "Publications vs. grant size",
                     _fig_pubs_vs_grant_scatter(df),
                     notes="Log-log scatter of publication count vs. grant size, with OLS trend (black).")
    _add_image_slide(prs, "Publication output by region",
                     _fig_publication_output_by_region(df),
                     notes="Regional ranking by total publications.")

    # ── Methodology ────────────────────────────────────────────────────────
    _add_section_divider(prs, "7.  Methodology",
                        notes="How the data was assembled and classified.")
    method_text = (
        "• Project registry: hand-curated Excel (UKRI_Projects_Partially_Cleaned.xlsx).\n"
        "• Enrichment: 16,128 per-project JSON files from UKRI Gateway-to-Research.\n"
        "    – `organisationRoles` → org–project participation index.\n"
        "    – `researchTopics` + `researchSubjects` → topic co-occurrence index.\n"
        "    – `personRoles` → investigator collaboration index.\n"
        "• Sustainability classification: keyword matching of title + abstract against "
        "eight pre-defined themes.\n"
        "• Community detection: greedy modularity on weighted graphs.\n"
        "• Plotly figures exported via kaleido for byte-identical match with the dashboard.\n"
        "• Network figures rendered with matplotlib (same code path as the dashboard)."
    )
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.4), Inches(12.5), Inches(6.7))
    tb.text_frame.word_wrap = True
    tb.text_frame.text = "Methodology"
    for p in tb.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(22); r.font.bold = True
            r.font.color.rgb = RGBColor(0x0B, 0x3D, 0x91)
    p2 = tb.text_frame.add_paragraph()
    p2.text = method_text
    for r in p2.runs:
        r.font.size = Pt(13); r.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    slide.notes_slide.notes_text_frame.text = (
        "Cite this slide for methodology questions. Full code in data_processor.py "
        "and the network/figure construction lives in app.py + pptx_export.py."
    )

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
